import os
import io
import faiss
import numpy as np
import google.generativeai as genai
import PyPDF2
from datetime import datetime
from typing import List
from dotenv import load_dotenv

load_dotenv()

# --- CONFIG ---
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
EMBEDDING_MODEL = "text-embedding-004"
GENERATION_MODEL = "gemini-2.5-flash"
dimension = 768  # embedding vector size

# --- Gemini Setup ---
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# --- FAISS Global State ---
index = None
chunks_store = []
embeddings_store = []
is_initialized = False

# Add conversation context tracking
conversation_context = {
    "current_document": None,
    "document_summary": "",
    "recent_questions": [],
    "last_context_used": ""
}

# -------------------- INIT --------------------
def initialize_faiss():
    global index, is_initialized
    index = faiss.IndexFlatL2(dimension)
    is_initialized = True

def is_index_initialized():
    """Check if FAISS index is initialized"""
    return is_initialized

def get_embeddings_count():
    """Get count of stored embeddings"""
    return len(embeddings_store)

def get_chunks_count():
    """Get count of stored chunks"""
    return len(chunks_store)

def get_current_timestamp():
    return datetime.now().isoformat()

# -------------------- CONVERSATION CONTEXT --------------------
def get_conversation_context():
    """Get current conversation context"""
    return conversation_context

def update_conversation_context(document_name=None, document_summary=None, question=None):
    """Update conversation context with new information"""
    global conversation_context
    
    if document_name:
        conversation_context["current_document"] = document_name
    if document_summary:
        conversation_context["document_summary"] = document_summary
    if question:
        conversation_context["recent_questions"].append(question)
        # Keep only last 5 questions for context
        conversation_context["recent_questions"] = conversation_context["recent_questions"][-5:]
    
    print(f"🔄 Conversation context updated: {conversation_context}")

def clear_conversation_context():
    """Clear conversation context when new document is uploaded"""
    global conversation_context
    conversation_context = {
        "current_document": None,
        "document_summary": "",
        "recent_questions": [],
        "last_context_used": ""
    }
    print("🧹 Conversation context cleared")

# -------------------- UNIVERSAL TEXT EXTRACTION --------------------
async def universal_text_extraction(file) -> str:
    """
    Universal text extraction pipeline that handles multiple file types
    Uses appropriate extraction method based on detected file type
    """
    try:
        # Detect file type
        file_info = detect_file_type(file)
        detected_type = file_info["type"]
        confidence = file_info["confidence"]
        method = file_info["method"]
        
        print(f"🔍 File type detection: {detected_type} (confidence: {confidence}, method: {method})")
        print(f"   File: {file_info['filename']}, MIME: {file_info['mime_type']}, Ext: {file_info['extension']}")
        
        # Route to appropriate extraction method
        if detected_type == "pdf":
            return await extract_pdf_text(file)
        elif detected_type in ["png", "jpg", "jpeg", "gif", "bmp", "tiff", "image"]:
            return await extract_image_text(file)
        elif detected_type in ["doc", "docx"]:
            return await extract_word_text(file)
        elif detected_type in ["ppt", "pptx"]:
            return await extract_powerpoint_text(file)
        elif detected_type in ["xls", "xlsx"]:
            return await extract_excel_text(file)
        elif detected_type in ["txt", "md", "markdown"]:
            return await extract_text_file(file)
        elif detected_type in ["html", "htm"]:
            return await extract_html_text(file)
        elif detected_type == "css":
            return await extract_css_text(file)
        elif detected_type in ["csv", "json", "xml"]:
            return await extract_structured_text(file)
        else:
            # Fallback: try to decode as text
            return await extract_fallback_text(file, file_info)
            
    except Exception as e:
        print(f"❌ Error in universal text extraction: {e}")
        return f"Error extracting text from file: {str(e)}"

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    Split text into overlapping chunks for better RAG performance.
    
    Args:
        text: The text to chunk
        chunk_size: Maximum size of each chunk in characters
        overlap: Number of characters to overlap between chunks
    
    Returns:
        List of text chunks
    """
    if not text or len(text.strip()) == 0:
        return []
    
    text = text.strip()
    chunks = []
    
    # If text is shorter than chunk size, return as single chunk
    if len(text) <= chunk_size:
        return [text]
    
    # Split text into chunks with overlap
    start = 0
    while start < len(text):
        end = start + chunk_size
        
        # If this is not the last chunk, try to break at a sentence boundary
        if end < len(text):
            # Look for sentence endings within the last 100 characters of the chunk
            search_start = max(start + chunk_size - 100, start)
            search_end = min(end + 50, len(text))
            
            # Find the last sentence ending in this range
            sentence_endings = ['.', '!', '?', '\n\n']
            best_break = end
            
            for ending in sentence_endings:
                last_occurrence = text.rfind(ending, search_start, search_end)
                if last_occurrence > start and last_occurrence < end:
                    best_break = last_occurrence + 1
                    break
            
            end = best_break
        
        # Extract the chunk
        chunk = text[start:end].strip()
        if chunk:  # Only add non-empty chunks
            chunks.append(chunk)
        
        # Move to next chunk with overlap
        start = end - overlap
        if start >= len(text):
            break
    
    print(f"📄 Text chunked into {len(chunks)} chunks (chunk_size: {chunk_size}, overlap: {overlap})")
    return chunks

def detect_file_type(file) -> dict:
    """
    Universal file type detection using multiple methods
    Returns: {"type": "pdf", "confidence": "high", "method": "mime+extension"}
    """
    try:
        filename = file.filename or ""
        content_type = file.content_type or ""
        
        # Method 1: MIME type detection
        mime_type = None
        if content_type:
            mime_type = content_type.lower()
        
        # Method 2: Extension-based detection
        extension = os.path.splitext(filename.lower())[1] if filename else ""
        
        # Method 3: Magic number detection (for binary files)
        magic_numbers = {
            b'%PDF': 'pdf',
            b'\x89PNG': 'png',
            b'\xff\xd8\xff': 'jpg',
            b'GIF8': 'gif',
            b'PK\x03\x04': 'zip',  # DOCX, PPTX, XLSX are ZIP-based
            b'\xd0\xcf\x11\xe0': 'doc',  # Old DOC format
        }
        
        # Read first few bytes for magic number detection
        file_content = file.file.read(8)
        file.file.seek(0)  # Reset file pointer
        
        detected_type = None
        confidence = "low"
        method = "unknown"
        
        # Check magic numbers first (highest confidence)
        for magic, file_type in magic_numbers.items():
            if file_content.startswith(magic):
                detected_type = file_type
                confidence = "high"
                method = "magic_number"
                break
        
        # Check MIME type
        if not detected_type and mime_type:
            if mime_type.startswith('image/'):
                detected_type = mime_type.split('/')[1]
                confidence = "high"
                method = "mime_type"
            elif mime_type in ['application/pdf', 'application/x-pdf']:
                detected_type = 'pdf'
                confidence = "high"
                method = "mime_type"
            elif mime_type in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
                detected_type = 'docx' if 'openxml' in mime_type else 'doc'
                confidence = "high"
                method = "mime_type"
            elif mime_type in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
                detected_type = 'pptx' if 'openxml' in mime_type else 'ppt'
                confidence = "high"
                method = "mime_type"
            elif mime_type in ['application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
                detected_type = 'xlsx' if 'openxml' in mime_type else 'xls'
                confidence = "high"
                method = "mime_type"
            elif mime_type in ['text/html', 'text/css', 'text/plain']:
                detected_type = mime_type.split('/')[1]
                confidence = "high"
                method = "mime_type"
        
        # Check file extension
        if not detected_type and extension:
            extension_map = {
                '.pdf': 'pdf',
                '.png': 'png', '.jpg': 'jpg', '.jpeg': 'jpeg', '.gif': 'gif', '.bmp': 'bmp', '.tiff': 'tiff',
                '.doc': 'doc', '.docx': 'docx',
                '.ppt': 'ppt', '.pptx': 'pptx',
                '.xls': 'xls', '.xlsx': 'xlsx',
                '.txt': 'txt', '.md': 'markdown',
                '.html': 'html', '.htm': 'html', '.css': 'css',
                '.csv': 'csv', '.json': 'json', '.xml': 'xml'
            }
            if extension in extension_map:
                detected_type = extension_map[extension]
                confidence = "medium" if method == "unknown" else confidence
                method = "extension" if method == "unknown" else method
        
        # Fallback detection
        if not detected_type:
            if 'pdf' in filename.lower():
                detected_type = 'pdf'
                confidence = "medium"
                method = "filename_pattern"
            elif any(img_ext in filename.lower() for img_ext in ['.png', '.jpg', '.jpeg', '.gif']):
                detected_type = 'image'
                confidence = "medium"
                method = "filename_pattern"
            elif any(doc_ext in filename.lower() for doc_ext in ['.doc', '.docx']):
                detected_type = 'docx' if '.docx' in filename.lower() else 'doc'
                confidence = "medium"
                method = "filename_pattern"
        
        return {
            "type": detected_type or "unknown",
            "confidence": confidence,
            "method": method,
            "mime_type": mime_type,
            "extension": extension,
            "filename": filename
        }
        
    except Exception as e:
        print(f"❌ Error in file type detection: {e}")
        return {
            "type": "unknown",
            "confidence": "low",
            "method": "error",
            "mime_type": content_type,
            "extension": "",
            "filename": filename
        }

# -------------------- ENHANCED EXTRACTION FUNCTIONS --------------------
async def extract_pdf_text(file) -> str:
    """Extract text from PDF files using multiple methods"""
    try:
        content = await file.read()
        pdf_file = io.BytesIO(content)
        
        # Try PyPDF2 first
        try:
            reader = PyPDF2.PdfReader(pdf_file)
            text = ""
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text += f"Page {page_num + 1}:\n{page_text}\n\n"
            
            if text.strip():
                print(f"✅ PDF text extracted with PyPDF2: {len(text)} characters")
                return text
        except Exception as e:
            print(f"⚠️ PyPDF2 failed: {e}")
            pdf_file.seek(0)  # Reset for next method
        
        # Try pdfplumber as fallback
        try:
            import pdfplumber
            with pdfplumber.open(pdf_file) as pdf:
                text = ""
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text += f"Page {page_num + 1}:\n{page_text}\n\n"
                
                if text.strip():
                    print(f"✅ PDF text extracted with pdfplumber: {len(text)} characters")
                    return text
        except ImportError:
            print("⚠️ pdfplumber not available")
        except Exception as e:
            print(f"⚠️ pdfplumber failed: {e}")
        
        # If both methods fail, return error
        return "PDF text extraction failed with all available methods. Please ensure the PDF contains extractable text."
        
    except Exception as e:
        print(f"❌ Error extracting PDF text: {e}")
        return f"Failed to extract PDF text: {str(e)}"

async def extract_image_text(file) -> str:
    """Extract text from images using OCR"""
    try:
        content = await file.read()
        
        # Try pytesseract first
        try:
            import pytesseract
            from PIL import Image
            
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
            
            if text.strip():
                print(f"✅ Image text extracted with pytesseract: {len(text)} characters")
                return f"Image: {file.filename}\n\nExtracted Text:\n{text}"
            else:
                print(f"⚠️ No text found in image: {file.filename}")
                return f"Image: {file.filename}\n\nNo text could be extracted from this image. It may contain only graphics, diagrams, or non-text content."
                
        except ImportError:
            print("⚠️ pytesseract not available")
            return f"Image: {file.filename}\n\nImage file detected but OCR library (pytesseract) not available. Please install it for text extraction."
        except Exception as e:
            print(f"⚠️ pytesseract failed: {e}")
            return f"Image: {file.filename}\n\nOCR processing failed: {str(e)}"
    
    except Exception as e:
        print(f"❌ Error extracting image text: {e}")
        return f"Failed to extract image text: {str(e)}"

async def extract_word_text(file) -> str:
    """Extract text from Word documents"""
    try:
        content = await file.read()
        
        if file.filename.lower().endswith('.docx'):
            # Process DOCX files
            try:
                from docx import Document
                doc = Document(io.BytesIO(content))
                text = ""
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text += paragraph.text + "\n"
                
                print(f"✅ DOCX text extracted: {len(text)} characters")
                return text
                
            except ImportError:
                return f"Word document: {file.filename}\n\nWord processing library (python-docx) not available. Please install it."
            except Exception as e:
                print(f"⚠️ DOCX processing failed: {e}")
                return f"Failed to process DOCX: {str(e)}"
                
        elif file.filename.lower().endswith('.doc'):
            # Process DOC files
            try:
                import docx2txt
                text = docx2txt.process(io.BytesIO(content))
                print(f"✅ DOC text extracted: {len(text)} characters")
                return text
                
            except ImportError:
                return f"Word document: {file.filename}\n\nWord processing library (docx2txt) not available. Please install it."
            except Exception as e:
                print(f"⚠️ DOC processing failed: {e}")
                return f"Failed to process DOC: {str(e)}"
        
        return f"Unsupported Word document format: {file.filename}"
        
    except Exception as e:
        print(f"❌ Error extracting Word text: {e}")
        return f"Failed to extract Word document text: {str(e)}"

async def extract_powerpoint_text(file) -> str:
    """Extract text from PowerPoint presentations"""
    try:
        content = await file.read()
        
        try:
            from pptx import Presentation
            
            prs = Presentation(io.BytesIO(content))
            text = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text += f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                text += "\n"
            
            if text.strip():
                print(f"✅ PowerPoint text extracted: {len(text)} characters")
                return text
            else:
                return f"PowerPoint: {file.filename}\n\nNo text content found in this presentation."
                
        except ImportError:
            return f"PowerPoint: {file.filename}\n\nPowerPoint processing library (python-pptx) not available. Please install it."
        except Exception as e:
            print(f"⚠️ PowerPoint processing failed: {e}")
            return f"Failed to process PowerPoint: {str(e)}"
        
    except Exception as e:
        print(f"❌ Error extracting PowerPoint text: {e}")
        return f"Failed to extract PowerPoint text: {str(e)}"

async def extract_excel_text(file) -> str:
    """Extract text from Excel spreadsheets"""
    try:
        content = await file.read()
        
        try:
            import pandas as pd
            
            # Try to read with pandas
            excel_file = io.BytesIO(content)
            
            # Read all sheets
            excel_data = pd.read_excel(excel_file, sheet_name=None)
            
            text = ""
            for sheet_name, df in excel_data.items():
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            
            if text.strip():
                print(f"✅ Excel text extracted: {len(text)} characters")
                return text
            else:
                return f"Excel: {file.filename}\n\nNo data found in this spreadsheet."
                
        except ImportError:
            return f"Excel: {file.filename}\n\nExcel processing library (pandas) not available. Please install it."
        except Exception as e:
            print(f"⚠️ Excel processing failed: {e}")
            return f"Failed to process Excel: {str(e)}"
        
    except Exception as e:
        print(f"❌ Error extracting Excel text: {e}")
        return f"Failed to extract Excel text: {str(e)}"

async def extract_text_file(file) -> str:
    """Extract text from plain text files"""
    try:
        content = await file.read()
        
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                text = content.decode(encoding)
                print(f"✅ Text file extracted with {encoding} encoding: {len(text)} characters")
                return text
            except UnicodeDecodeError:
                continue
        
        # If all encodings fail, try with error handling
        try:
            text = content.decode('utf-8', errors='replace')
            print(f"⚠️ Text file extracted with replacement characters: {len(text)} characters")
            return text
        except Exception:
            return f"Text file: {file.filename}\n\nFailed to decode text content with any supported encoding."
        
    except Exception as e:
        print(f"❌ Error extracting text file: {e}")
        return f"Failed to extract text file: {str(e)}"

async def extract_html_text(file) -> str:
    """Extract text from HTML files"""
    try:
        content = await file.read()
        
        try:
            from bs4 import BeautifulSoup
            
            # Parse HTML and extract text
            soup = BeautifulSoup(content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Get text and clean it up
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            if text.strip():
                print(f"✅ HTML text extracted: {len(text)} characters")
                return text
            else:
                return f"HTML: {file.filename}\n\nNo text content found in this HTML file."
                
        except ImportError:
            return f"HTML: {file.filename}\n\nHTML processing library (beautifulsoup4) not available. Please install it."
        except Exception as e:
            print(f"⚠️ HTML processing failed: {e}")
            return f"Failed to process HTML: {str(e)}"
        
    except Exception as e:
        print(f"❌ Error extracting HTML text: {e}")
        return f"Failed to extract HTML text: {str(e)}"

async def extract_css_text(file) -> str:
    """Extract text from CSS files"""
    try:
        content = await file.read()
        
        try:
            text = content.decode('utf-8')
            
            # Remove CSS comments and clean up
            import re
            # Remove /* ... */ comments
            text = re.sub(r'/\*.*?\*/', '', text, flags=re.DOTALL)
            # Remove // comments
            text = re.sub(r'//.*$', '', text, flags=re.MULTILINE)
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text).strip()
            
            if text.strip():
                print(f"✅ CSS text extracted: {len(text)} characters")
                return f"CSS File: {file.filename}\n\nExtracted CSS Rules:\n{text}"
            else:
                return f"CSS: {file.filename}\n\nNo CSS rules found in this file."
                
        except Exception as e:
            print(f"⚠️ CSS processing failed: {e}")
            return f"Failed to process CSS: {str(e)}"
        
    except Exception as e:
        print(f"❌ Error extracting CSS text: {e}")
        return f"Failed to extract CSS text: {str(e)}"

async def extract_structured_text(file) -> str:
    """Extract text from structured data files (CSV, JSON, XML)"""
    try:
        content = await file.read()
        filename = file.filename.lower()
        
        if filename.endswith('.csv'):
            try:
                import pandas as pd
                df = pd.read_csv(io.BytesIO(content))
                text = f"CSV File: {file.filename}\n\nData:\n{df.to_string(index=False)}"
                print(f"✅ CSV text extracted: {len(text)} characters")
                return text
            except ImportError:
                return f"CSV: {file.filename}\n\nCSV processing library (pandas) not available. Please install it."
            except Exception as e:
                return f"Failed to process CSV: {str(e)}"
                
        elif filename.endswith('.json'):
            try:
                import json
                data = json.loads(content.decode('utf-8'))
                text = f"JSON File: {file.filename}\n\nData:\n{json.dumps(data, indent=2)}"
                print(f"✅ JSON text extracted: {len(text)} characters")
                return text
            except Exception as e:
                return f"Failed to process JSON: {str(e)}"
                
        elif filename.endswith('.xml'):
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content, 'xml')
                text = f"XML File: {file.filename}\n\nData:\n{soup.prettify()}"
                print(f"✅ XML text extracted: {len(text)} characters")
                return text
            except ImportError:
                return f"XML: {file.filename}\n\nXML processing library (beautifulsoup4) not available. Please install it."
            except Exception as e:
                return f"Failed to process XML: {str(e)}"
        
        return f"Unsupported structured file format: {file.filename}"
        
    except Exception as e:
        print(f"❌ Error extracting structured text: {e}")
        return f"Failed to extract structured text: {str(e)}"

async def extract_fallback_text(file, file_info: dict) -> str:
    """Fallback text extraction for unknown file types"""
    try:
        content = await file.read()
        
        # Try to decode as text with various encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                text = content.decode(encoding)
                if text.strip():
                    print(f"✅ Fallback text extraction with {encoding}: {len(text)} characters")
                    return f"File: {file_info['filename']}\nType: {file_info['mime_type']}\n\nExtracted Text:\n{text}"
            except UnicodeDecodeError:
                continue
        
        # If text decoding fails, return file info
        return f"File: {file_info['filename']}\nType: {file_info['mime_type']}\nSize: {len(content)} bytes\n\nThis file type is not specifically supported for text extraction, but it has been uploaded successfully."
        
    except Exception as e:
        print(f"❌ Error in fallback text extraction: {e}")
        return f"File: {file_info['filename']}\nType: {file_info['mime_type']}\n\nFile uploaded but text extraction failed: {str(e)}"

# -------------------- FILE PROCESSING --------------------
async def process_document(file) -> str:
    """Detects file type and extracts text from multiple formats."""
    content_type = (file.content_type or "").lower()
    filename = file.filename.lower()
    data = await file.read()

    if "pdf" in content_type or filename.endswith(".pdf"):
        return await extract_pdf_text(file)

    elif filename.endswith((".docx", ".doc")):
        return await extract_word_text(file)

    elif filename.endswith((".pptx", ".ppt")):
        return await extract_powerpoint_text(file)

    elif filename.endswith((".xlsx", ".xls")):
        return await extract_excel_text(file)

    elif filename.endswith((".html", ".htm", ".css", ".txt")):
        return await extract_text_file(file)

    elif content_type.startswith("image/") or filename.endswith((".png", ".jpg", ".jpeg")):
        return await extract_image_text(file)

    else:
        return f"Unsupported file type: {filename}"


# -------------------- RAG CORE --------------------
async def create_embeddings(texts: List[str]) -> List[dict]:
    """Creates embeddings for text chunks using Gemini API."""
    global embeddings_store, chunks_store
    
    try:
        if not is_initialized:
            initialize_faiss()
        
        # Clear previous conversation context when new documents are uploaded
        clear_conversation_context()
        
        # Store chunks
        chunks_store = texts
        
        # Create a summary of the document for context
        if texts:
            document_summary = texts[0][:500] + "..." if len(texts[0]) > 500 else texts[0]
            update_conversation_context(document_summary=document_summary)
        
        if GEMINI_API_KEY:
            # Use Gemini embedding API
            embeddings = await embed_texts_gemini(texts)
        else:
            # Fallback to random embeddings (for testing)
            print("⚠️ No Gemini API key found, using random embeddings")
            embeddings = create_random_embeddings(texts)
        
        # Store embeddings
        embeddings_store = embeddings
        
        # Add to FAISS index
        if index is not None:
            embedding_vectors = np.array([emb["embedding"] for emb in embeddings], dtype="float32")
            index.add(embedding_vectors)
            print(f"💾 {len(embeddings)} embeddings added to FAISS index")
        
        return embeddings
    
    except Exception as e:
        print(f"❌ Error creating embeddings: {e}")
        raise Exception(f"Failed to create embeddings: {str(e)}")

async def embed_texts_gemini(texts: List[str]) -> List[dict]:
    """Create embeddings using Gemini API"""
    try:
        model = EMBEDDING_MODEL
        embeddings = []
        
        for i, text in enumerate(texts):
            try:
                result = genai.embed_content(model=model, content=text)
                # API may return dict-like with 'embedding'
                embedding_vector = result.get("embedding") if isinstance(result, dict) else getattr(result, "embedding", None)
                if embedding_vector is None:
                    raise ValueError("No embedding returned from Gemini")
                
                embeddings.append({
                    "id": i,
                    "text": text,
                    "embedding": embedding_vector,
                    "metadata": {
                        "chunk_index": i,
                        "timestamp": get_current_timestamp(),
                        "model": model
                    }
                })
                
            except Exception as e:
                print(f"⚠️ Error embedding chunk {i}: {e}")
                # Fallback to random embedding
                embeddings.append({
                    "id": i,
                    "text": text,
                    "embedding": np.random.rand(dimension).tolist(),
                    "metadata": {
                        "chunk_index": i,
                        "timestamp": get_current_timestamp(),
                        "model": "fallback-random"
                    }
                })
        
        return embeddings
    
    except Exception as e:
        print(f"❌ Error with Gemini embeddings: {e}")
        raise Exception(f"Gemini embedding failed: {str(e)}")

def create_random_embeddings(texts: List[str]) -> List[dict]:
    """Create random embeddings as fallback"""
    embeddings = []
    
    for i, text in enumerate(texts):
        embeddings.append({
            "id": i,
            "text": text,
            "embedding": np.random.rand(dimension).tolist(),
            "metadata": {
                "chunk_index": i,
                "timestamp": get_current_timestamp(),
                "model": "fallback-random"
            }
        })
    
    return embeddings

async def answer_question(question: str) -> tuple[str, float]:
    """Answer a question using RAG with stored embeddings and Gemini's general knowledge. Returns (answer, confidence_score)"""
    try:
        # Update conversation context with the current question
        update_conversation_context(question=question)
        
        # Debug logging to see what's stored
        print(f"🔍 Debug - chunks_store length: {len(chunks_store)}")
        print(f"🔍 Debug - embeddings_store length: {len(embeddings_store)}")
        print(f"🔍 Debug - index initialized: {is_initialized}")
        if index is not None:
            print(f"🔍 Debug - FAISS index total vectors: {index.ntotal}")
        
        # Initialize context and confidence
        document_context = ""
        confidence_score = 0.0
        
        # If documents are uploaded, search for relevant content
        if chunks_store and embeddings_store and is_initialized:
            print(f"📚 Documents available, searching for relevant content...")
            # Create question embedding
            if GEMINI_API_KEY:
                question_embedding = await embed_texts_gemini([question])
                q_vector = np.array([question_embedding[0]["embedding"]], dtype="float32")
            else:
                # Fallback to random embedding
                q_vector = np.random.rand(1, dimension).astype("float32")
            
            # Search for relevant chunks
            if index is not None and index.ntotal > 0:
                try:
                    D, I = index.search(q_vector, k=min(3, len(chunks_store)))
                    relevant_chunks = [chunks_store[i] for i in I[0] if i < len(chunks_store)]
                    document_context = "\n\n".join(relevant_chunks)
                    
                    # Calculate confidence based on similarity scores and context relevance
                    if len(D[0]) > 0:
                        # Normalize distance scores (lower distance = higher similarity = higher confidence)
                        # Convert to Python floats to avoid NumPy array comparison issues
                        distances = [float(d) for d in D[0]]
                        max_distance = max(distances) if distances else 1.0
                        min_distance = min(distances) if distances else 0.0
                        if max_distance > min_distance:
                            # Convert distance to similarity (0-1 scale)
                            avg_similarity = 1.0 - (sum(distances) / len(distances) / max_distance)
                            # Boost confidence if we have substantial context
                            context_boost = min(len(document_context) / 1000.0, 0.3)  # Max 0.3 boost for long context
                            confidence_score = min(avg_similarity + context_boost, 1.0)
                        else:
                            confidence_score = 0.5  # Default confidence for equal distances
                    else:
                        confidence_score = 0.3  # Low confidence if no search results
                    
                    print(f"📚 Found {len(relevant_chunks)} relevant document chunks")
                    print(f"📚 Context preview: {document_context[:200]}...")
                    print(f"📊 Calculated confidence score: {confidence_score:.3f}")
                except Exception as search_error:
                    print(f"⚠️ FAISS search failed: {search_error}")
                    # Fallback: use first few chunks without similarity scoring
                    relevant_chunks = chunks_store[:min(3, len(chunks_store))]
                    document_context = "\n\n".join(relevant_chunks)
                    confidence_score = 0.2  # Low confidence due to search failure
                    print(f"📚 Using fallback chunks: {len(relevant_chunks)} chunks")
            else:
                print(f"⚠️ FAISS index is empty or not properly initialized")
                confidence_score = 0.1  # Very low confidence
        else:
            print(f"⚠️ No documents uploaded yet or system not initialized")
            print(f"   chunks_store: {len(chunks_store)}, embeddings_store: {len(embeddings_store)}, is_initialized: {is_initialized}")
            confidence_score = 0.0  # No confidence without documents
        
        # Generate answer using PDF-first approach
        if GEMINI_API_KEY:
            clean_context = sanitize_document_context(document_context)
            print(f"🧹 Cleaned context length: {len(clean_context)}")
            answer = await generate_pdf_first_response(question, clean_context)
        else:
            answer = generate_fallback_response(question, document_context)
            confidence_score = max(confidence_score, 0.2)  # Minimum confidence for fallback
        
        return answer, confidence_score
    
    except Exception as e:
        print(f"❌ Error answering question: {e}")
        error_response = f"Sorry, I encountered an error while processing your question: {str(e)}"
        return error_response, 0.0

async def generate_pdf_first_response(question: str, document_context: str = "") -> str:
    """Generate response using PDF-first approach: check PDF context first, then use general knowledge if needed"""
    try:
        model = genai.GenerativeModel(GENERATION_MODEL)
        
        # Get conversation context for better understanding
        context = get_conversation_context()
        
        # Enforce conversational tone, no unwanted bullets
        system_instruction = (
            "You are a helpful AI tutor. "
            "Answer in a friendly, concise conversational style. "
            "Do not use bullet points or numbered lists unless the user explicitly asks for them. "
            "Break into short sentences for clarity. "
            "Write as if you're having a natural conversation with a student. "
            "IMPORTANT: When the user asks about 'this', 'that', or 'it', always refer to the uploaded document content."
        )
        
        if document_context:
            # Use new prompt strategy: check PDF first, then general knowledge
            prompt = f"""{system_instruction}

You are analyzing an uploaded document. The user's question may reference "this", "that", or "it" - these refer to the document content below.

First, check if the answer is in the provided PDF context.  
If it is, answer strictly based on that context.  
If it is not, answer using your own general knowledge while making it clear that this is outside the uploaded PDF material.

=== DOCUMENT CONTEXT START ===
{document_context}
=== DOCUMENT CONTEXT END ===

Current Document: {context.get('current_document', 'Uploaded document')}

Question: {question}

Instructions:
- If the question asks about "this", "that", or "it", refer to the document content above
- If the question is vague like "explain this", analyze the document content and explain what it's about
- Provide a clear, student-friendly explanation in conversational language
- If using outside knowledge, say: "From my own knowledge: ..."

Answer:"""
        else:
            # No PDF content available, use general knowledge
            prompt = f"""{system_instruction}

The user has not uploaded a PDF yet, so I'll answer using my general knowledge.

Question: {question}

Answer in clear, conversational language:"""
        
        response = model.generate_content(prompt)
        
        # Format the response for better readability
        formatted_response = format_response_for_user(response.text, question, document_context)
        return formatted_response
    
    except Exception as e:
        print(f"❌ Error with Gemini API: {e}")
        return generate_fallback_response(question, document_context)

def sanitize_document_context(raw_context: str) -> str:
    """Remove placeholder or non-informative contexts so the model doesn't get confused."""
    if not raw_context:
        return ""
    lowered = raw_context.lower()
    noisy_markers = [
        "processing libraries not available",
        "ocr library (pytesseract) not available",
        "no text could be extracted",
        "file uploaded but text extraction failed",
        "word document detected but",
        "library not available",
        "please install it",
        "failed to process",
        "failed to extract"
    ]
    if any(marker in lowered for marker in noisy_markers):
        return ""
    if len(raw_context.strip()) < 30:
        return ""
    return raw_context

def format_response_for_user(raw_response: str, question: str, document_context: str = "") -> str:
    """Format the raw AI response into a clean, conversational format"""
    try:
        # Clean up the response
        cleaned = raw_response.strip()
        
        # If response contains the question, remove it for cleaner output
        if question.lower() in cleaned.lower():
            # Try to find where the actual answer starts
            lines = cleaned.split('\n')
            answer_lines = []
            found_answer = False
            
            for line in lines:
                line = line.strip()
                if line and not line.lower().startswith('question:') and not line.lower().startswith('answer:'):
                    found_answer = True
                if found_answer:
                    answer_lines.append(line)
            
            if answer_lines:
                cleaned = '\n'.join(answer_lines)
        
        # Remove any remaining prompt artifacts
        artifacts_to_remove = [
            "You are a helpful AI tutor.",
            "Answer in a friendly, concise conversational style.",
            "Do not use bullet points or numbered lists unless the user explicitly asks for them.",
            "Break into short sentences for clarity.",
            "Write as if you're having a natural conversation with a student.",
            "First, check if the answer is in the provided PDF context.",
            "If it is, answer strictly based on that context.",
            "If it is not, answer using your own general knowledge while making it clear that this is outside the uploaded PDF material.",
            "=== DOCUMENT CONTEXT START ===",
            "=== DOCUMENT CONTEXT END ===",
            "Provide a clear, student-friendly explanation in conversational language.",
            "If using outside knowledge, say: \"From my own knowledge: ...\"",
            "The user has not uploaded a PDF yet, so I'll answer using my general knowledge.",
            "Answer in clear, conversational language:"
        ]
        
        for artifact in artifacts_to_remove:
            cleaned = cleaned.replace(artifact, "").replace(artifact.strip(), "")
        
        # Clean up extra whitespace and formatting
        cleaned = '\n'.join(line.strip() for line in cleaned.split('\n') if line.strip())
        
        # If the response is too short or seems incomplete, return the original
        if len(cleaned) < 20:
            return raw_response
        
        return cleaned
        
    except Exception as e:
        print(f"⚠️ Error formatting response: {e}")
        # Return original response if formatting fails
        return raw_response

async def expand_query(query: str) -> str:
    """Expand/rewrite vague user queries into more explicit, context-rich queries"""
    try:
        if not GEMINI_API_KEY:
            print("⚠️ No Gemini API key available, returning original query")
            return query
        
        # Get current conversation context
        context = get_conversation_context()
        
        # Check if this is a contextual reference that needs expansion
        contextual_indicators = [
            "this", "that", "it", "the document", "the pdf", "the file", "the image",
            "explain this", "what is this", "tell me about this", "describe this",
            "what's this about", "what is this about", "explain that", "what is that"
        ]
        
        query_lower = query.lower().strip()
        is_contextual = any(indicator in query_lower for indicator in contextual_indicators)
        
        if is_contextual and context["current_document"]:
            # This is a contextual reference to the uploaded document
            expanded_query = f"Based on the uploaded document '{context['current_document']}', {query}"
            print(f"🔍 Contextual query detected and expanded: '{query}' -> '{expanded_query}'")
            return expanded_query
        
        # Create a prompt to expand the query
        expansion_prompt = f"""You are a helpful AI assistant that rewrites vague or incomplete user questions to make them more explicit and searchable.

Original Question: "{query}"

Current Context: {context['current_document'] or 'No document uploaded'}

Instructions:
1. If the question is already clear and specific, return it as-is
2. If the question is vague, incomplete, or ambiguous, rewrite it to be more explicit
3. If the question references "this", "that", or "it" and there's a document context, make it clear what "this" refers to
4. Add context and specificity where helpful
5. Keep the rewritten question concise and focused
6. Maintain the original intent and meaning
7. Use natural, conversational language

Rewritten Question:"""
        
        # Call Gemini to expand the query
        expanded = await call_gemini(expansion_prompt)
        
        # Clean up the response
        expanded = expanded.strip()
        
        # If Gemini returned something useful, use it; otherwise fall back to original
        if expanded and len(expanded) > 10 and expanded.lower() != "error":
            print(f"🔍 Query expanded: '{query}' -> '{expanded}'")
            return expanded
        else:
            print(f"⚠️ Query expansion failed, using original: '{query}'")
            return query
    
    except Exception as e:
        print(f"❌ Error expanding query: {e}")
        return query

async def call_gemini(prompt: str) -> str:
    """Call Gemini API directly with a custom prompt"""
    try:
        if not GEMINI_API_KEY:
            return "Gemini API not configured"
        
        model = genai.GenerativeModel(GENERATION_MODEL)
        response = model.generate_content(prompt)
        return response.text
    
    except Exception as e:
        print(f"❌ Error calling Gemini directly: {e}")
        return f"Error calling Gemini: {str(e)}"

def generate_fallback_response(question: str, context: str = "") -> str:
    """Generate fallback response when Gemini is not available"""
    if context:
        return f"""Based on the uploaded document, here's what I found:

Context: {context[:200]}...

Question: {question}

Note: This is a fallback response. For better AI-powered answers, please configure the Gemini API key."""
    else:
        return f"""Question: {question}

Note: This is a fallback response. For better AI-powered answers, please configure the Gemini API key. I can only provide basic responses without the AI model.""" 

# Initialize FAISS when module is imported
if not is_initialized:
    initialize_faiss()
